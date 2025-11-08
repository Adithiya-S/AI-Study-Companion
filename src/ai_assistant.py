"""
AI Assistant Module - Gemini AI integration for study materials
Supports both internet-based and uploaded materials-based queries
Compatible with Python 3.8-3.12
"""

import os
import json
from pathlib import Path
from typing import Dict, List, Optional, Any
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed

try:
    from dotenv import load_dotenv
    DOTENV_AVAILABLE = True
except ImportError:
    DOTENV_AVAILABLE = False
    
try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False
    print("Warning: google-generativeai not installed. AI features will be disabled.")

# Document parsing libraries
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    
try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False


class AIAssistant:
    """Manages AI-powered study assistance using Gemini API."""
    
    def __init__(self, data_dir: str = "data"):
        self.data_dir = Path(data_dir)
        self.materials_dir = self.data_dir / "materials"
        self.uploaded_files_dir = self.materials_dir / "uploaded_files"
        self.chat_history_file = self.materials_dir / "chat_history.json"
        
        # Ensure directories exist
        self.uploaded_files_dir.mkdir(parents=True, exist_ok=True)
        
        # AI configuration
        self.api_key = None
        self.model = None
        self.chat_session = None
        self.chat_history = []
        # Debounced persistence settings
        self._last_history_save_ts = 0.0
        self._history_save_min_interval = 2.0  # seconds
        
        # Load API key from .env file
        self.load_api_key_from_env()
        
        # Initialize Gemini if available
        if GEMINI_AVAILABLE and self.api_key:
            self.initialize_gemini()
        
        # Load chat history
        self.load_chat_history()
        
    def load_api_key_from_env(self):
        """Load Gemini API key from .env file."""
        try:
            # Load environment variables from .env file
            if DOTENV_AVAILABLE:
                # Try to find .env in project root
                env_path = Path(__file__).parent.parent / '.env'
                if env_path.exists():
                    load_dotenv(env_path)
            
            # Get API key from environment variable
            self.api_key = os.getenv('GEMINI_API_KEY')
            
            if not self.api_key:
                print("Warning: GEMINI_API_KEY not found in environment variables.")
                print("Please create a .env file with your API key.")
                
        except Exception as e:
            print(f"Error loading API key from environment: {e}")
            
    def initialize_gemini(self):
        """Initialize Gemini AI model."""
        try:
            if not self.api_key:
                return False
                
            genai.configure(api_key=self.api_key)
            
            # Use Gemini 2.0 Flash model (stable, latest generation)
            # Alternatives: 'gemini-1.5-flash', 'gemini-1.5-pro'
            self.model = genai.GenerativeModel('gemini-2.0-flash')
            
            # Start a new chat session
            self.chat_session = self.model.start_chat(history=[])
            
            return True
            
        except Exception as e:
            print(f"Error initializing Gemini: {e}")
            return False
            
    def is_configured(self) -> bool:
        """Check if AI assistant is properly configured."""
        return GEMINI_AVAILABLE and self.api_key is not None and self.model is not None
        
    def upload_study_material(self, file_path: str, title: str = "", 
                             description: str = "") -> Dict[str, Any]:
        """
        Upload and process a study material file.
        Supports: txt, pdf, md, docx (text-based files)
        """
        try:
            source_path = Path(file_path)
            
            if not source_path.exists():
                return {"success": False, "error": "File not found"}
                
            # Generate unique filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            dest_filename = f"{timestamp}_{source_path.name}"
            dest_path = self.uploaded_files_dir / dest_filename
            
            # Read file content
            content = self.extract_text_from_file(source_path)
            
            if not content:
                return {"success": False, "error": "Could not extract text from file"}
                
            # Save to uploaded files directory with metadata
            metadata = {
                "original_name": source_path.name,
                "title": title or source_path.stem,
                "description": description,
                "upload_date": datetime.now().isoformat(),
                "file_type": source_path.suffix,
                "content": content[:1000],  # Preview
                "full_path": str(dest_path),
                "word_count": len(content.split())
            }
            
            # Copy file to uploaded directory
            with open(dest_path, 'w', encoding='utf-8') as f:
                f.write(content)
                
            # Save metadata
            metadata_path = dest_path.with_suffix('.meta.json')
            with open(metadata_path, 'w', encoding='utf-8') as f:
                json.dump(metadata, f, indent=2, ensure_ascii=False)
                
            return {
                "success": True,
                "file_id": dest_filename,
                "metadata": metadata
            }
            
        except Exception as e:
            return {"success": False, "error": str(e)}
            
    def extract_text_from_file(self, file_path: Path) -> str:
        """Extract text content from various file formats."""
        try:
            suffix = file_path.suffix.lower()
            
            # Plain text files
            if suffix in ['.txt', '.md', '.py', '.java', '.cpp', '.js', '.html', '.css', '.json', '.xml', '.csv']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            
            # PDF files
            elif suffix == '.pdf' and PDF_AVAILABLE:
                return self._extract_from_pdf(file_path)
            
            # Word documents
            elif suffix in ['.docx', '.doc'] and DOCX_AVAILABLE:
                return self._extract_from_docx(file_path)
            
            # PowerPoint presentations
            elif suffix in ['.pptx', '.ppt'] and PPTX_AVAILABLE:
                return self._extract_from_pptx(file_path)
            
            # Excel spreadsheets
            elif suffix in ['.xlsx', '.xls'] and XLSX_AVAILABLE:
                return self._extract_from_xlsx(file_path)
            
            else:
                # Check if library is available for the format
                if suffix == '.pdf' and not PDF_AVAILABLE:
                    return "PDF support not available. Please install PyPDF2."
                elif suffix in ['.docx', '.doc'] and not DOCX_AVAILABLE:
                    return "Word document support not available. Please install python-docx."
                elif suffix in ['.pptx', '.ppt'] and not PPTX_AVAILABLE:
                    return "PowerPoint support not available. Please install python-pptx."
                elif suffix in ['.xlsx', '.xls'] and not XLSX_AVAILABLE:
                    return "Excel support not available. Please install openpyxl."
                else:
                    return f"Unsupported file type: {suffix}"
                
        except Exception as e:
            print(f"Error extracting text: {e}")
            return f"Error extracting text: {str(e)}"
    
    def _extract_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file."""
        try:
            text_parts = []
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                num_pages = len(pdf_reader.pages)
                
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text.strip():
                        text_parts.append(f"--- Page {page_num + 1} ---\n{text}\n")
            
            return "\n".join(text_parts) if text_parts else "No text could be extracted from PDF."
        except Exception as e:
            return f"Error extracting PDF: {str(e)}"
    
    def _extract_from_docx(self, file_path: Path) -> str:
        """Extract text from Word document."""
        try:
            doc = Document(file_path)
            text_parts = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        text_parts.append(row_text)
            
            return "\n".join(text_parts) if text_parts else "No text could be extracted from Word document."
        except Exception as e:
            return f"Error extracting Word document: {str(e)}"
    
    def _extract_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # More than just the header
                    text_parts.append("\n".join(slide_text))
            
            return "\n\n".join(text_parts) if text_parts else "No text could be extracted from PowerPoint."
        except Exception as e:
            return f"Error extracting PowerPoint: {str(e)}"
    
    def _extract_from_xlsx(self, file_path: Path) -> str:
        """Extract text from Excel spreadsheet."""
        try:
            wb = load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"--- Sheet: {sheet_name} ---"]
                
                for row in sheet.iter_rows(values_only=True):
                    # Filter out None values and convert to strings
                    row_values = [str(cell) for cell in row if cell is not None]
                    if row_values:
                        sheet_text.append(" | ".join(row_values))
                
                if len(sheet_text) > 1:  # More than just the header
                    text_parts.append("\n".join(sheet_text))
            
            return "\n\n".join(text_parts) if text_parts else "No data could be extracted from Excel file."
        except Exception as e:
            return f"Error extracting Excel file: {str(e)}"
            
    def get_uploaded_materials(self) -> List[Dict[str, Any]]:
        """Get list of all uploaded study materials."""
        materials = []
        
        try:
            for meta_file in self.uploaded_files_dir.glob("*.meta.json"):
                with open(meta_file, 'r', encoding='utf-8') as f:
                    metadata = json.load(f)
                    materials.append(metadata)
                    
            # Sort by upload date (newest first)
            materials.sort(key=lambda x: x.get("upload_date", ""), reverse=True)
            
        except Exception as e:
            print(f"Error getting uploaded materials: {e}")
            
        return materials
        
    def delete_uploaded_material(self, file_id: str) -> bool:
        """Delete an uploaded material and its metadata."""
        try:
            file_path = self.uploaded_files_dir / file_id
            meta_path = file_path.with_suffix('.meta.json')
            
            if file_path.exists():
                file_path.unlink()
            if meta_path.exists():
                meta_path.unlink()
                
            return True
            
        except Exception as e:
            print(f"Error deleting material: {e}")
            return False
            
    def ask_question(self, question: str, mode: str = "internet", 
                    context_files: Optional[List[str]] = None) -> Dict[str, Any]:
        """
        Ask a question to the AI assistant.
        
        Args:
            question: The question to ask
            mode: "internet" for web-based answers, "materials" for uploaded materials only
            context_files: List of file IDs to use as context (for materials mode)
            
        Returns:
            Dict with answer and metadata
        """
        if not self.is_configured():
            return {
                "success": False,
                "error": "AI Assistant not configured. Please add your Gemini API key."
            }
            
        try:
            if mode == "materials":
                return self.ask_from_materials(question, context_files)
            else:
                return self.ask_from_internet(question)
                
        except Exception as e:
            return {
                "success": False,
                "error": f"Error processing question: {str(e)}"
            }
            
    def ask_from_internet(self, question: str) -> Dict[str, Any]:
        """Ask a question with internet knowledge."""
        try:
            # Create prompt for internet-based query
            prompt = f"""You are a helpful study assistant. Answer the following question clearly and concisely:

Question: {question}

Please provide:
1. A clear, accurate answer
2. Key points or concepts
3. If relevant, suggestions for further study

Answer:"""

            response = self.chat_session.send_message(prompt)
            answer = response.text
            
            result = {
                "success": True,
                "answer": answer,
                "mode": "internet",
                "timestamp": datetime.now().isoformat(),
                "question": question
            }
            
            # Save to chat history
            self.add_to_history(result)
            
            return result
            
        except Exception as e:
            return {
                "success": False,
                "error": f"Error getting internet-based answer: {str(e)}"
            }
            
    def ask_from_materials(self, question: str, 
                          context_files: Optional[List[str]] = None) -> Dict[str, Any]:
        """Ask a question based only on uploaded materials."""
        try:
            # Get relevant materials content
            if context_files:
                # Use specific files
                materials_content = self.load_specific_materials(context_files)
            else:
                # Use all uploaded materials
                materials_content = self.load_all_materials()
                
            if not materials_content:
                return {
                    "success": True,
                    "answer": "No information found in the uploaded materials. Please upload study materials first or switch to Internet mode.",
                    "mode": "materials",
                    "found_info": False,
                    "timestamp": datetime.now().isoformat(),
                    "question": question
                }
                
            # Create prompt with materials context
            prompt = f"""You are a study assistant. Answer the question ONLY using the information from the uploaded study materials below. If the information is not in the materials, respond with "No information found in the uploaded materials."

Study Materials:
{materials_content}

Question: {question}

Instructions:
- Only use information from the study materials above
- If the answer is not in the materials, say "No information found in the uploaded materials."
- Be specific and cite relevant parts of the materials
- Keep the answer clear and focused

Answer:"""

            response = self.chat_session.send_message(prompt)
            answer = response.text
            
            # Check if answer indicates no information found
            found_info = "no information found" not in answer.lower()
            
            result = {
                "success": True,
                "answer": answer,
                "mode": "materials",
                "found_info": found_info,
                "timestamp": datetime.now().isoformat(),
                "question": question,
                "materials_used": len(context_files) if context_files else "all"
            }
            
            # Save to chat history
            self.add_to_history(result)
            
            return result
            
        except Exception as e:
            return {
                "success": False,
                "error": f"Error getting materials-based answer: {str(e)}"
            }
            
    def load_specific_materials(self, file_ids: List[str]) -> str:
        """Load content from specific material files."""
        content_parts = []
        
        for file_id in file_ids:
            try:
                file_path = self.uploaded_files_dir / file_id
                if file_path.exists():
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                        content_parts.append(f"=== {file_id} ===\n{content}\n")
            except Exception as e:
                print(f"Error loading material {file_id}: {e}")
                
        return "\n".join(content_parts)
        
    def load_all_materials(self) -> str:
        """Load content from all uploaded materials."""
        content_parts = []
        try:
            files = [p for p in self.uploaded_files_dir.iterdir() if p.suffix != '.json']
            if not files:
                return ""
            # Parallelize file reads (I/O bound)
            def _read(p: Path) -> str:
                try:
                    with open(p, 'r', encoding='utf-8') as f:
                        return f"=== {p.name} ===\n{f.read()}\n"
                except Exception:
                    return ""
            with ThreadPoolExecutor(max_workers=min(8, len(files))) as ex:
                futures = [ex.submit(_read, p) for p in files]
                for fut in as_completed(futures):
                    part = fut.result()
                    if part:
                        content_parts.append(part)
        except Exception as e:
            print(f"Error loading all materials: {e}")
        return "\n".join(content_parts)
        
    def add_to_history(self, interaction: Dict[str, Any]):
        """Add interaction to chat history."""
        self.chat_history.append(interaction)
        self.save_chat_history()
        
    def load_chat_history(self):
        """Load chat history from file."""
        try:
            if self.chat_history_file.exists():
                with open(self.chat_history_file, 'r', encoding='utf-8') as f:
                    self.chat_history = json.load(f)
        except Exception as e:
            print(f"Error loading chat history: {e}")
            self.chat_history = []
            
    def save_chat_history(self):
        """Save chat history to file."""
        try:
            # Debounce disk writes
            now = datetime.now().timestamp()
            if now - self._last_history_save_ts < self._history_save_min_interval:
                return
            # Keep only last 100 interactions
            if len(self.chat_history) > 100:
                self.chat_history = self.chat_history[-100:]
            with open(self.chat_history_file, 'w', encoding='utf-8') as f:
                json.dump(self.chat_history, f, indent=2, ensure_ascii=False)
            self._last_history_save_ts = now
        except Exception as e:
            print(f"Error saving chat history: {e}")
            
    def get_chat_history(self, limit: int = 20) -> List[Dict[str, Any]]:
        """Get recent chat history."""
        return self.chat_history[-limit:] if self.chat_history else []
        
    def clear_chat_history(self):
        """Clear all chat history."""
        self.chat_history = []
        self.save_chat_history()
        
        # Reset chat session
        if self.model:
            self.chat_session = self.model.start_chat(history=[])
    
    def get_chat_sessions(self) -> List[Dict[str, Any]]:
        """Get list of chat sessions grouped by timestamp."""
        sessions = []
        if not self.chat_history:
            return sessions
        
        # Group chats by date
        current_session = []
        last_timestamp = None
        
        for item in self.chat_history:
            timestamp = item.get("timestamp", "")
            if timestamp:
                try:
                    dt = datetime.fromisoformat(timestamp)
                    # If more than 30 minutes gap, start new session
                    if last_timestamp and (dt - last_timestamp).total_seconds() > 1800:
                        if current_session:
                            sessions.append({
                                "start_time": current_session[0].get("timestamp"),
                                "end_time": current_session[-1].get("timestamp"),
                                "message_count": len(current_session),
                                "preview": current_session[0].get("question", "")[:50] + "...",
                                "messages": current_session.copy()
                            })
                        current_session = []
                    
                    current_session.append(item)
                    last_timestamp = dt
                except:
                    pass
        
        # Add last session
        if current_session:
            sessions.append({
                "start_time": current_session[0].get("timestamp"),
                "end_time": current_session[-1].get("timestamp"),
                "message_count": len(current_session),
                "preview": current_session[0].get("question", "")[:50] + "...",
                "messages": current_session.copy()
            })
        
        # Reverse to show most recent first
        sessions.reverse()
        return sessions
    
    def start_new_chat(self):
        """Start a fresh chat session without loading history."""
        if self.model:
            self.chat_session = self.model.start_chat(history=[])
        return True
